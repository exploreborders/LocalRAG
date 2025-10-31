from langchain_community.document_loaders import TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
import os
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

def load_documents(data_dir="data"):
    """
    Load documents from the specified directory.

    Supports multiple file formats: .txt, .pdf, .docx, .pptx, .xlsx

    Args:
        data_dir (str): Path to the directory containing documents

    Returns:
        list: List of Document objects with page_content and metadata
    """
    documents = []
    for file in os.listdir(data_dir):
        file_path = os.path.join(data_dir, file)
        if file.endswith(".txt"):
            loader = TextLoader(file_path)
            documents.extend(loader.load())
        elif file.endswith(".pdf"):
            documents.extend(load_pdf(file_path))
        elif file.endswith(".docx"):
            documents.extend(load_docx(file_path))
        elif file.endswith(".pptx"):
            documents.extend(load_pptx(file_path))
        elif file.endswith(".xlsx"):
            documents.extend(load_xlsx(file_path))
        # Add more formats as needed
    return documents

def load_pdf(file_path):
    """
    Load a PDF file and extract text content.

    Args:
        file_path (str): Path to the PDF file

    Returns:
        list: List containing a single Document object with extracted text
    """
    documents = []
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        documents.append(Document(page_content=text, metadata={"source": file_path}))
    return documents

def load_docx(file_path):
    """
    Load a DOCX file and extract text content.

    Args:
        file_path (str): Path to the DOCX file

    Returns:
        list: List containing a single Document object with extracted text
    """
    documents = []
    doc = DocxDocument(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    documents.append(Document(page_content=text, metadata={"source": file_path}))
    return documents

def load_pptx(file_path):
    """
    Load a PPTX file and extract text content from slides.

    Args:
        file_path (str): Path to the PPTX file

    Returns:
        list: List containing a single Document object with extracted text
    """
    documents = []
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    documents.append(Document(page_content=text, metadata={"source": file_path}))
    return documents

def load_xlsx(file_path):
    """
    Load an XLSX file and extract text content from all sheets.

    Args:
        file_path (str): Path to the XLSX file

    Returns:
        list: List containing a single Document object with extracted text
    """
    documents = []
    wb = openpyxl.load_workbook(file_path)
    text = ""
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
    documents.append(Document(page_content=text, metadata={"source": file_path}))
    return documents

def split_documents(documents, chunk_size=1000, chunk_overlap=200):
    """
    Split documents into smaller chunks for better retrieval performance.

    Args:
        documents (list): List of Document objects to split
        chunk_size (int): Maximum size of each chunk in characters
        chunk_overlap (int): Number of characters to overlap between chunks

    Returns:
        list: List of Document objects representing text chunks
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )
    chunks = text_splitter.split_documents(documents)
    return chunks

if __name__ == "__main__":
    docs = load_documents()
    print(f"Loaded {len(docs)} documents")
    chunks = split_documents(docs)
    print(f"Split into {len(chunks)} chunks")
    for i, chunk in enumerate(chunks[:3]):  # Show first 3 chunks
        print(f"Chunk {i+1}: {chunk.page_content[:200]}...")